"""
create_ppt.py — Deepfake Detection Project
Creates a 10-slide Review-1 presentation using the provided template.
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os
import sys

TEMPLATE_PATH = sys.argv[1] if len(sys.argv) > 1 else r"C:\Users\HP VICTUS\Downloads\Missed_Lead_Detector_Presentation.pptx"
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "..", "outputs", "Deepfake_Detection_Review1.pptx")

# ── Load template ─────────────────────────────────────────
prs = Presentation(TEMPLATE_PATH)

# Remove all existing slides (we'll build fresh using the template's layouts/theme)
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId is None:
        rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


# ── Helper Functions ─────────────────────────────────────
def add_title_slide(prs, title_text, subtitle_text=""):
    """Add a title slide matching the template's first slide style."""
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:  # Title placeholder
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = title_text
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.name = "Cambria"
            run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
        elif ph.placeholder_format.idx == 1:  # Subtitle placeholder
            tf = ph.text_frame
            tf.clear()
            if subtitle_text:
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = subtitle_text
                run.font.size = Pt(14)
                run.font.name = "Times New Roman"
                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return slide


def add_content_slide(prs, title_text, bullets, sub_bullets=None):
    """Add a content slide with title and bullet points."""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:  # Title
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title_text
            run.font.size = Pt(24)
            run.font.bold = True
            run.font.name = "Cambria"
            run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
        elif ph.placeholder_format.idx == 1:  # Content
            tf = ph.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.space_after = Pt(6)
                p.space_before = Pt(3)
                run = p.add_run()
                run.text = bullet
                run.font.size = Pt(13)
                run.font.name = "Times New Roman"
                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                # Add sub-bullets if any
                if sub_bullets and i in sub_bullets:
                    for sb in sub_bullets[i]:
                        sp = tf.add_paragraph()
                        sp.level = 1
                        sp.space_after = Pt(3)
                        srun = sp.add_run()
                        srun.text = sb
                        srun.font.size = Pt(11)
                        srun.font.name = "Times New Roman"
                        srun.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    return slide





def add_thank_you_slide(prs):
    """Add a thank you slide."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = "Thank You!"
            run.font.size = Pt(44)
            run.font.bold = True
            run.font.name = "Cambria"
            run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
        elif ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = "AI-Based Deepfake Detection and Verification System\nUsing Deep Learning"
            run.font.size = Pt(16)
            run.font.name = "Times New Roman"
            run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    return slide


# ═══════════════════════════════════════════════════════════
#  SLIDE 1: TITLE SLIDE
# ═══════════════════════════════════════════════════════════
add_title_slide(
    prs,
    title_text="AI-Based Deepfake Detection\nand Verification System\nUsing Deep Learning",
    subtitle_text="Review-1: Problem Identification & Title Justification\n\nTeam Members:\n[Your Name] — [Regd No.]\n[Member 2 Name] — [Regd No.]\n[Member 3 Name] — [Regd No.]\n[Member 4 Name] — [Regd No.]\n\n[Your College] | [Your Dept] | Batch [Year]"
)


# ═══════════════════════════════════════════════════════════
#  SLIDE 2: RELEVANCE OF THE PROBLEM
# ═══════════════════════════════════════════════════════════
add_content_slide(
    prs,
    title_text="Relevance of the Problem",
    bullets=[
        "The rapid growth of AI has enabled creation of highly realistic fake images and videos (deepfakes)",
        "Deepfakes are used for fake news, identity theft, financial fraud, and political manipulation",
        "Social media misinformation — manipulated content spreads faster than authentic content",
        "Traditional manual verification methods are unreliable and time-consuming",
        "As deepfakes become increasingly difficult to identify with the human eye, automated detection is critical",
        "Growing need for accurate systems to distinguish genuine vs. manipulated digital media",
    ],
    sub_bullets={
        0: ["Industries affected: Banking, Social Media, Journalism, Law Enforcement, Politics"],
        2: ["Cost of deepfake fraud exceeded $25 billion globally in 2023 (FTC Report)"],
        4: ["Humans can only detect ~50% of deepfakes in controlled studies (Nature, 2022)"],
    }
)


# ═══════════════════════════════════════════════════════════
#  SLIDE 3: PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════
add_content_slide(
    prs,
    title_text="Problem Statement",
    bullets=[
        "Deepfake technology has advanced significantly — fake videos and images appear highly realistic",
        "Manual verification is unreliable, slow, and cannot scale with the volume of online content",
        "No universally deployed lightweight system exists for real-time deepfake detection",
        "Existing detection models require significant computational resources",
        "Therefore, an intelligent system is required to automatically analyze digital media and identify whether it is authentic or manipulated",
    ],
    sub_bullets={
        0: ["GANs, Autoencoders, and Diffusion Models enable high-quality face synthesis"],
        1: ["Social media platforms process billions of images/videos daily"],
        3: ["Most state-of-the-art models run on GPUs — not suitable for edge/mobile deployment"],
    }
)


# ═══════════════════════════════════════════════════════════
#  SLIDE 4: JUSTIFICATION OF PROJECT TITLE
# ═══════════════════════════════════════════════════════════
add_content_slide(
    prs,
    title_text="Justification of Project Title",
    bullets=[
        '"AI-Based" — Uses Artificial Intelligence and Deep Learning algorithms',
        '"Deepfake Detection" — Identifies manipulated facial images and videos',
        '"Verification System" — Provides binary verification (Real / Fake) with confidence scores',
        '"Using Deep Learning" — Employs CNNs, Transfer Learning, and advanced neural architectures',
        "The system contributes to cybersecurity and digital media authenticity",
    ],
    sub_bullets={
        0: ["CNNs, Transfer Learning (VGG, ResNet, EfficientNet), and ensemble methods"],
        1: ["Analyzes facial inconsistencies: lighting, texture, blending artifacts, blinking patterns"],
        2: ["Confidence score indicates probability of manipulation (0% to 100%)"],
        3: ["Pre-trained models fine-tuned on deepfake datasets for high accuracy with less training data"],
    }
)


# ═══════════════════════════════════════════════════════════
#  SLIDE 5: OBJECTIVES OF THE PROJECT
# ═══════════════════════════════════════════════════════════
add_content_slide(
    prs,
    title_text="Objectives of the Project",
    bullets=[
        "Main: Develop an intelligent deep learning-based system for accurate deepfake detection",
        "Detect manipulated images and videos using facial feature analysis",
        "Analyze facial inconsistencies (lighting, texture, blending, blinking patterns)",
        "Achieve high detection accuracy with optimized model architecture",
        "Provide real-time verification capability for practical deployment",
        "Develop a user-friendly web interface for easy upload and verification",
        "Compare different deep learning models and select the best performer",
    ],
    sub_bullets={
        0: ["Target: AUC > 0.98 and Accuracy > 95% on benchmark datasets"],
        2: ["Face extraction using MTCNN / dlib before classification"],
        5: ["Upload image/video → Instant prediction with confidence score"],
        6: ["CNN, VGG16, ResNet50, EfficientNet, and custom architectures"],
    }
)


# ═══════════════════════════════════════════════════════════
#  SLIDE 6: SCOPE OF THE PROJECT
# ═══════════════════════════════════════════════════════════
add_content_slide(
    prs,
    title_text="Scope of the Project",
    bullets=[
        "Included in Scope:",
        "Deepfake image detection (face-level manipulation)",
        "Deepfake video detection (frame-by-frame + temporal analysis)",
        "Face extraction and preprocessing pipeline",
        "Real-time prediction with confidence score generation",
        "Performance evaluation (Accuracy, AUC, F1-Score, Confusion Matrix)",
        "Future Scope:",
        "Voice / audio deepfake detection (speech synthesis)",
        "Social media API integration for live content scanning",
        "Browser extension for real-time image verification",
        "Mobile application deployment (Android / iOS)",
        "Multimodal detection combining video + audio + text analysis",
    ],
    sub_bullets={
        1: ["Supports JPEG, PNG, and BMP image formats"],
        2: ["Extracts keyframes for efficient video analysis"],
        4: ["5-fold cross-validation, ROC-AUC curves, precision-recall analysis"],
        8: ["Lightweight TFLite / ONNX models for on-device inference"],
    }
)


# ═══════════════════════════════════════════════════════════
#  SLIDE 7: BASIC CONCEPTS RELATED TO THE PROJECT
# ═══════════════════════════════════════════════════════════
add_content_slide(
    prs,
    title_text="Basic Concepts Related to the Project",
    bullets=[
        "Artificial Intelligence (AI) — Simulation of human intelligence using machines",
        "Machine Learning (ML) — Algorithms that learn patterns from data without explicit programming",
        "Deep Learning (DL) — Advanced ML using multi-layered neural networks",
        "Computer Vision — AI field enabling computers to understand images and videos",
        "Convolutional Neural Network (CNN) — Deep learning architecture for image classification",
        "Deepfake — AI-generated fake media that appears authentic to humans",
        "Transfer Learning — Using pre-trained models to improve performance and reduce training time",
    ],
    sub_bullets={
        0: ["Rule-based systems → Statistical ML → Deep Learning → Generative AI"],
        4: ["Convolution layers extract spatial features; pooling reduces dimensionality"],
        5: ["Generative Adversarial Networks (GANs) and Autoencoders are primary generation methods"],
        6: ["Models like VGG16, ResNet50, EfficientNet pre-trained on ImageNet, fine-tuned for deepfakes"],
    }
)


# ═══════════════════════════════════════════════════════════
#  SLIDE 8: LITERATURE SURVEY
# ═══════════════════════════════════════════════════════════
add_content_slide(
    prs,
    title_text="Literature Survey",
    bullets=[
        '[1] "FaceForensics++" (Rössler et al., 2019) — Introduced benchmark dataset for deepfake detection',
        '[2] "Celeb-DF" (Li et al., 2020) — Large-scale challenging dataset with improved deepfake quality',
        '[3] "DeepFake Detection Using CNN-Based Models" — Achieved good classification accuracy',
        '[4] "Vision Transformer for Deepfake Detection" — Captures global image features effectively',
        '[5] "Hybrid CNN and Transformer Models" — Combines local and global feature extraction',
    ],
    sub_bullets={
        0: ["Contribution: Standard benchmark | Limitation: Performance drops on compressed videos"],
        1: ["Contribution: High-quality synthesis | Limitation: Struggles with highly realistic deepfakes"],
        2: ["Contribution: Good accuracy | Limitation: High computational complexity"],
        3: ["Contribution: Global features | Limitation: Requires large computational resources"],
        4: ["Contribution: Best of both | Limitation: Increased model complexity"],
    }
)


# ═══════════════════════════════════════════════════════════
#  SLIDE 9: RESEARCH GAP & PROPOSED SOLUTION
# ═══════════════════════════════════════════════════════════
add_content_slide(
    prs,
    title_text="Research Gap & Proposed Solution",
    bullets=[
        "High computational cost → Lightweight CNN architecture (fewer parameters)",
        "Slow detection speed → Real-time inference pipeline",
        "Large model size → Optimized model via pruning & quantization",
        "Limited deployment → Web-based Flask/Streamlit interface",
        "Poor generalization → Transfer learning on pre-trained models",
        "Proposed Research Gap:",
        "Most existing deepfake detection systems achieve high accuracy but require significant",
        "computational resources and are not suitable for real-time applications. This project",
        "proposes a lightweight and efficient deep learning framework capable of providing",
        "accurate real-time deepfake detection while maintaining lower computational complexity.",
    ],
    sub_bullets={
        5: ["This is the core thesis statement of our project"],
    }
)


# ═══════════════════════════════════════════════════════════
#  SLIDE 10: PROJECT PLANNING & TIMELINE
# ═══════════════════════════════════════════════════════════
add_content_slide(
    prs,
    title_text="Project Planning & Timeline (16 Weeks)",
    bullets=[
        "Week 1–2: Literature Survey & Research Gap Identification",
        "Week 3–4: Dataset Collection & Analysis (FaceForensics++, Celeb-DF, DFDC)",
        "Week 5–6: Data Preprocessing & Face Extraction Pipeline",
        "Week 7–8: Baseline CNN Model Development & Initial Training",
        "Week 9–10: Model Improvement — Transfer Learning, Hyperparameter Tuning",
        "Week 11–12: Real-Time System Development (Web Interface + API)",
        "Week 13: Testing, Evaluation & Performance Comparison",
        "Week 14–15: Documentation & Research Paper Drafting",
        "Week 16: Final Review & Submission",
        "Milestones: M1 — Literature Survey  |  M2 — Dataset Prepared  |  M3 — Baseline CNN",
        "M4 — Improved DL Model  |  M5 — Web Interface  |  M6 — Evaluation  |  M7 — Paper Drafted  |  M8 — Final Submission",
        "Task Allocation: Member 1 — Literature & Docs  |  Member 2 — Dataset & Preprocessing",
        "Member 3 — Model Development & Training  |  Member 4 — Frontend & Deployment",
    ]
)


# ── Save ──────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
prs.save(OUTPUT_PATH)
print(f"Presentation saved to: {OUTPUT_PATH}")
print(f"Total slides: {len(prs.slides)}")
