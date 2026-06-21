"""
generate_presentation.py — Missed-Lead Detector
Generates updated PBL Review 1 PPT with full project content + architecture diagram.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = os.path.dirname(__file__)
OUT_DIR = os.path.join(BASE, "..")
PPTX_OUT = os.path.join(OUT_DIR, "Missed_Lead_Detector_PBL_Review1.pptx")

# ── Color Palette (Aurora Borealis Dark) ──
BG_DARK    = RGBColor(0x0A, 0x0A, 0x12)
BG_CARD    = RGBColor(0x14, 0x18, 0x22)
BLUE       = RGBColor(0x38, 0xBD, 0xF8)
GREEN      = RGBColor(0x34, 0xD3, 0x99)
RED        = RGBColor(0xF8, 0x71, 0x71)
AMBER      = RGBColor(0xFB, 0xBF, 0x24)
PURPLE     = RGBColor(0xA7, 0x8B, 0xFA)
WHITE      = RGBColor(0xE6, 0xED, 0xF3)
GRAY       = RGBColor(0x7D, 0x85, 0x90)
DARK_GRAY  = RGBColor(0x48, 0x4F, 0x58)
ACCENT     = RGBColor(0x25, 0x63, 0xEB)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf

def add_bullet_points(slide, left, top, width, height, items, font_size=16, color=WHITE, bullet_color=BLUE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸ {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return tf

def add_card(slide, left, top, width, height, title, items, title_color=BLUE, bg_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + 0.2), Inches(width - 0.6), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.color.rgb = title_color
    p.font.bold = True

    add_bullet_points(slide, left + 0.3, top + 0.65, width - 0.6, height - 0.85, items, font_size=13, color=GRAY)

# ══════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide1, BG_DARK)

add_textbox(slide1, 1, 1.5, 11, 1, "Missed-Lead Detector",
            font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide1, 1, 2.6, 11, 0.8,
            "Machine Learning-Based Missed-Lead Detection and Automated\nFollow-Up System for Customer Retention in Sales Pipelines",
            font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide1, 1, 3.8, 11, 0.5, "AD4V71 — Machine Learning Operations (MLOps)",
            font_size=16, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# Divider line
shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.5), Inches(5), Inches(0.02))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_textbox(slide1, 1, 4.8, 11, 0.5, "Presented By",
            font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide1, 1, 5.3, 11, 1,
            "Team Member 1 Name  •  Regd No.  •  Section\nTeam Member 2 Name  •  Regd No.  •  Section",
            font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide1, 1, 6.5, 11, 0.4, "CIT Chennai — Batch 2025-27",
            font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 2: Problem Identification
# ══════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, BG_DARK)

add_textbox(slide2, 0.5, 0.3, 12, 0.7, "Problem Identification & Title Justification",
            font_size=30, color=WHITE, bold=True)

shape = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(3), Inches(0.03))
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()

add_card(slide2, 0.5, 1.3, 5.8, 2.8, "🔍 Relevance of the Problem", [
    "Customers message businesses daily about price, demo, or availability — these are 'leads'",
    "Busy staff sometimes miss replying; after ~24 hours, unanswered customers move to competitors",
    "This 'missed lead' problem is common for SMBs with limited staff",
    "No real-time way to know which messages were not replied to timely",
    "A missed lead is rarely noticed until the sale is already lost",
], title_color=RED)

add_card(slide2, 6.8, 1.3, 5.8, 2.8, "🎯 Scope of the Project", [
    "Title: ML-Based Missed-Lead Detection & Automated Follow-Up System",
    "ML identifies leads missed by staff, then auto-sends follow-up or alerts staff",
    "Built and tested on self-created sample + Kaggle real-world data",
    "Real business data integration planned for next phase",
    "8 classification models tested including Deep Learning & Ensembles",
], title_color=GREEN)

add_card(slide2, 0.5, 4.4, 12.1, 2.8, "📋 Project Title Justification", [
    "Machine Learning identifies which customer messages did not receive a timely reply",
    "Automated Follow-Up Engine sends human-like replies (clients cannot detect automation)",
    "Sales Pipeline Retention — ensures no business opportunity is lost due to missed messages",
    "System includes: Gmail monitoring → ML scoring → Auto-reply → Dashboard → Notifications",
    "GitHub Actions schedules automatic scans every 10 minutes",
], title_color=AMBER)

# ══════════════════════════════════════════════════════════
# SLIDE 3: Objectives
# ══════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, BG_DARK)

add_textbox(slide3, 0.5, 0.3, 12, 0.7, "Objectives of the Project",
            font_size=30, color=WHITE, bold=True)

shape = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(3), Inches(0.03))
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()

add_card(slide3, 0.5, 1.3, 5.8, 2.5, "🎯 Primary Objectives", [
    "Build a system that analyzes customer messages and identifies untimely replies",
    "Automatically send follow-up emails to missed leads",
    "Alert staff members with reminders until they respond",
    "Ensure no business opportunity is lost due to missed messages",
    "Goal: 100% lead coverage with automated detection",
], title_color=BLUE)

add_card(slide3, 6.8, 1.3, 5.8, 2.5, "📊 Technical Objectives", [
    "Train 8 ML models and compare performance (AUC, F1, CV scores)",
    "Implement Optuna hyperparameter tuning for XGBoost",
    "Build PyTorch deep learning model with residual connections",
    "Create intent-aware smart reply engine (8 intent categories)",
    "Deploy Streamlit dashboard for real-time monitoring",
], title_color=PURPLE)

add_card(slide3, 0.5, 4.1, 12.1, 3.1, "🚀 Current Scope & Next Phase", [
    "Currently: Self-created synthetic data (500 leads) + Kaggle datasets (9,240 + 4,000 leads = 13,740 total)",
    "Currently: Text-based messages — email, chat, WhatsApp, phone inquiry channels",
    "Currently: Gmail IMAP integration for real inbox monitoring",
    "Next Phase: Connect with real business CRM data sources",
    "Next Phase: Multi-platform integration (WhatsApp Business API, website chat widgets)",
    "Next Phase: Real-time WebSocket dashboard updates",
], title_color=AMBER)

# ══════════════════════════════════════════════════════════
# SLIDE 4: Basic Concepts & ML Models
# ══════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, BG_DARK)

add_textbox(slide4, 0.5, 0.3, 12, 0.7, "Basic Concepts & Machine Learning Models",
            font_size=30, color=WHITE, bold=True)

shape = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(3), Inches(0.03))
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()

# Key Features
add_card(slide4, 0.5, 1.3, 4.0, 2.3, "🔧 Key Features", [
    "Response Gap: time between message and reply",
    "Intent Score: customer interest level from words",
    "Channel Encoding: email, chat, WhatsApp, phone",
    "Business Hours Flag: 9AM–6PM detection",
    "Gap Bucket: 0-6h, 6-12h, 12-24h, 24h+",
    "Message Length: text complexity proxy",
], title_color=BLUE)

# Traditional ML
add_card(slide4, 4.7, 1.3, 4.0, 2.3, "🤖 Traditional ML Models", [
    "Logistic Regression (Baseline)",
    "Naive Bayes (Probabilistic)",
    "Decision Tree (Rule-based)",
    "Random Forest (Bagging Ensemble)",
    "XGBoost (Gradient Boosting)",
    "Optuna-tuned XGBoost (Best: AUC 0.9794)",
], title_color=GREEN)

# Advanced Models
add_card(slide4, 8.9, 1.3, 3.9, 2.3, "🧠 Advanced Models", [
    "Voting Ensemble (RF+XGB+LR)",
    "Grand Ensemble (ML+DL soft-voting)",
    "Deep Learning: PyTorch Neural Network",
    "  - Residual blocks + BatchNorm",
    "  - He init + GELU activation",
    "  - Class-weighted BCE loss",
], title_color=PURPLE)

# Unsupervised
add_card(slide4, 0.5, 3.9, 4.0, 1.6, "📊 Unsupervised Learning", [
    "KMeans Clustering (3 segments)",
    "  - High-Intent-Missed",
    "  - Low-Intent",
    "  - Already-Converted",
], title_color=AMBER)

# Hyperparameter Tuning
add_card(slide4, 4.7, 3.9, 4.0, 1.6, "⚙️ Optuna Tuning", [
    "100 trials with TPE sampler",
    "5-fold stratified cross-validation",
    "Optimized: n_estimators, max_depth,",
    "  learning_rate, subsample, gamma",
], title_color=RED)

# Model Performance
add_card(slide4, 8.9, 3.9, 3.9, 1.6, "🏆 Model Performance", [
    "XGBoost Tuned: AUC 0.9794",
    "Random Forest: AUC 0.9725",
    "Grand Ensemble: AUC 0.9709",
    "Deep Learning: AUC 0.9613",
], title_color=GREEN)

# Intent Categories
add_card(slide4, 0.5, 5.7, 12.3, 1.6, "🎯 Smart Reply Engine — 8 Intent Categories", [
    "Pricing | Demo | Course | Placement | Complaint | Interest | Availability | Urgent",
    "Each intent has multiple template variations for natural variety — clients cannot tell replies are automated",
    "Intent detection uses keyword scoring with normalized confidence scores across all 8 categories",
], title_color=BLUE)

# ══════════════════════════════════════════════════════════
# SLIDE 5: System Architecture (with embedded diagram image)
# ══════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, BG_DARK)

add_textbox(slide5, 0.5, 0.3, 12, 0.7, "System Architecture & Pipeline",
            font_size=30, color=WHITE, bold=True)

shape = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(3), Inches(0.03))
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()

# Pipeline flow boxes
pipeline_steps = [
    ("📧 Gmail Inbox\n(IMAP Fetch)", BLUE),
    ("🔍 Newsletter\nFilter", GRAY),
    ("🤖 ML Scoring\n(8 Models)", GREEN),
    ("💬 Smart Reply\nEngine", PURPLE),
    ("📤 Auto-Reply\n(SMTP)", AMBER),
    ("🔔 Notifications\n(Email+Dashboard)", RED),
]

for i, (label, color) in enumerate(pipeline_steps):
    x = 0.5 + i * 2.1
    shape = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.4), Inches(1.8), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    txBox = slide5.shapes.add_textbox(Inches(x + 0.1), Inches(1.5), Inches(1.6), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Arrow between boxes
    if i < len(pipeline_steps) - 1:
        arrow = slide5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.85), Inches(1.7), Inches(0.2), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = DARK_GRAY
        arrow.line.fill.background()

# Architecture details (compact row)
add_card(slide5, 0.5, 2.8, 6.0, 1.8, "📧 Email Integration (IMAP/SMTP)", [
    "IMAP4_SSL fetches Gmail inbox with promotional filter",
    "Deduplicates via Message-ID tracking",
    "SMTP threaded replies with In-Reply-To headers",
], title_color=BLUE)

add_card(slide5, 6.8, 2.8, 6.0, 1.8, "🤖 ML Scoring Pipeline", [
    "Grand Ensemble: 50% ML + 50% DL (soft-voting)",
    "XGBoost tuned with Optuna (100 trials, TPE sampler)",
    "Feature engineering: 9 features per lead",
], title_color=GREEN)

# ── Embedded Architecture Diagram Image ──
ARCH_IMG = os.path.join(BASE, "..", "outputs", "system_architecture.png")
if os.path.exists(ARCH_IMG):
    # Place the architecture diagram image centered below the pipeline
    slide5.shapes.add_picture(ARCH_IMG, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.5))
else:
    # Fallback: show detail cards if image not found
    add_card(slide5, 0.5, 4.8, 6.0, 2.4, "📊 Streamlit Dashboard", [
        "Command Center: KPIs, pipeline controls, notifications",
        "Lead Explorer: Filterable table + smart reply simulator",
        "Model Analytics: 8-model AUC comparison, DL charts",
        "Auto-Replies Tracker: Follow-up status management",
    ], title_color=PURPLE)

    add_card(slide5, 6.8, 4.8, 6.0, 2.4, "⚙️ Infrastructure & Scheduling", [
        "GitHub Actions: Automatic 10-minute inbox scans",
        "Follow-Up Tracker: Flags leads >24h without human reply",
        "Overdue Escalation: Desktop popups + email alerts at 48h",
        "Streamlit Cloud: Auto-deploy on GitHub push",
    ], title_color=AMBER)

# ══════════════════════════════════════════════════════════
# SLIDE 6: Results & Performance
# ══════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, BG_DARK)

add_textbox(slide6, 0.5, 0.3, 12, 0.7, "Results & Model Performance",
            font_size=30, color=WHITE, bold=True)

shape = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(3), Inches(0.03))
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()

# Model comparison table
models = [
    ("XGBoost (Optuna-tuned)", "0.9794", "🏆 Best ML", GREEN),
    ("RandomForest", "0.9725", "", BLUE),
    ("Ensemble (RF+XGB+LR)", "0.9720", "", BLUE),
    ("Grand Ensemble (ML+DL)", "0.9709", "🎯 Overall Best", PURPLE),
    ("Deep Learning (PyTorch)", "0.9613", "", PURPLE),
    ("Decision Tree", "0.9627", "", BLUE),
    ("Logistic Regression", "0.9132", "Baseline", GRAY),
    ("Naive Bayes", "0.8604", "", GRAY),
]

# Table header
add_card(slide6, 0.5, 1.3, 8.0, 0.5, "", ["Model                                              AUC Score    Note"], title_color=BLUE, bg_color=RGBColor(0x1E, 0x24, 0x30))

for i, (name, auc, note, color) in enumerate(models):
    y = 1.9 + i * 0.55
    bg = BG_CARD if i % 2 == 0 else RGBColor(0x18, 0x1C, 0x26)
    add_card(slide6, 0.5, y, 8.0, 0.45, "",
             [f"{name:<40s} {auc:<12s} {note}"],
             title_color=color, bg_color=bg)

# Key metrics
add_card(slide6, 8.8, 1.3, 4.0, 1.5, "📊 Dataset", [
    "Total: 13,740 samples",
    "Synthetic: 500 leads",
    "Kaggle Lead Scoring: 9,240",
    "Kaggle Support Tickets: 4,000",
], title_color=BLUE)

add_card(slide6, 8.8, 3.0, 4.0, 1.5, "⚙️ Training Config", [
    "5-fold Stratified CV",
    "80/20 Train-Test Split",
    "StandardScaler normalization",
    "Early stopping (patience=20)",
], title_color=GREEN)

add_card(slide6, 8.8, 4.7, 4.0, 1.5, "🔧 Hyperparameters", [
    "XGBoost: 100 trials (Optuna)",
    "DL: 200 epochs max, batch=64",
    "Dropout: 0.3, LR: 1e-3",
    "Weight decay: 1e-4",
], title_color=PURPLE)

add_card(slide6, 8.8, 6.2, 4.0, 1.0, "💡 Key Insight", [
    "Grand Ensemble (ML+DL) achieves",
    "best overall performance by combining",
    "strengths of both paradigms",
], title_color=AMBER)

# ══════════════════════════════════════════════════════════
# SLIDE 7: Thank You
# ══════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7, BG_DARK)

add_textbox(slide7, 1, 2.0, 11, 1.5, "Thank You!",
            font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide7, 1, 3.5, 11, 0.8,
            "Missed-Lead Detector — AI-Powered Sales Command Center",
            font_size=20, color=BLUE, alignment=PP_ALIGN.CENTER)

shape = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.3), Inches(3), Inches(0.03))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_textbox(slide7, 1, 4.8, 11, 0.8,
            "CIT Chennai • Batch 2025-27\nAD4V71 — Machine Learning Operations (MLOps)",
            font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# ── Save ──
prs.save(PPTX_OUT)
print(f"[pptx] Saved updated presentation -> {PPTX_OUT}")
print(f"[pptx] {len(prs.slides)} slides generated")
